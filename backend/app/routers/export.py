"""Export router: generate Word (.docx) and PowerPoint (.pptx) files."""

from __future__ import annotations

import datetime
import importlib.resources
import io
import re

import requests
from docx.shared import Mm
from docxtpl import DocxTemplate, InlineImage
from fastapi import APIRouter, HTTPException, status
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

import highlight as hlt
from ..deps import require_session
from ..schemas import PptExportRequest, WordExportRequest
from ..session import Session

router = APIRouter(prefix="/export", tags=["export"])

_USER_AGENT = (
    "PAIGE/1.0 (https://github.com/crvernon/highlight; chris.vernon@pnnl.gov)"
)
_DOCX_MIME = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


def _build_base_filename(citation: str | None) -> str:
    """Derive an export base filename from a Chicago-style citation."""
    today_str = datetime.date.today().strftime("%d%b%Y").lower()
    if not citation:
        return f"ber-highlight_{today_str}"

    last_name = "unknown"
    year = "YYYY"
    journal_abbrev = "journal"

    try:
        first_author_part = citation.split(",")[0].strip()
        if " " in first_author_part and not first_author_part.isupper():
            last_name = first_author_part.split(" ")[0].strip().lower()
        else:
            last_name = first_author_part.lower()
        last_name = re.sub(r"[^a-z\-]", "", last_name) or "unknown"
    except Exception:  # noqa: BLE001
        last_name = "unknown"

    try:
        match = re.search(r"[\.\(;]\s*(\d{4})[\.\);]", citation)
        if match:
            year = match.group(1)
        else:
            match = re.search(r"\b(\d{4})\b", citation)
            if match:
                year = match.group(1)
    except Exception:  # noqa: BLE001
        year = "YYYY"

    try:
        match = re.search(r'["”]\s*\.\s*([^,.:]+?)\s*\d+[:\(]', citation, re.IGNORECASE)
        if match:
            journal_abbrev = re.sub(r"[\s\.,]", "", match.group(1).strip()).lower()
        elif "applied energy" in citation.lower():
            journal_abbrev = "appliedenergy"
        journal_abbrev = journal_abbrev or "journal"
    except Exception:  # noqa: BLE001
        journal_abbrev = "journal"

    return f"{last_name}_etal_{year}_{journal_abbrev}_ber-highlight_{today_str}"


def _build_photo_site_name(req: WordExportRequest) -> str:
    """Construct the attribution string used in the Word template."""
    if not req.selected_image:
        return "Wikimedia Commons"
    info = req.selected_image
    parts: list[str] = []
    if info.artist_plain and info.artist_plain != "Unknown Artist":
        parts.append(info.artist_plain)
    if info.license:
        license_part = info.license
        if info.license_url:
            license_part += f" <{info.license_url}>"
        parts.append(license_part)
    parts.append("via Wikimedia Commons")
    return ", ".join(filter(None, parts)).replace("<", "(").replace(">", ")")


def _fetch_image_bytes(url: str | None) -> io.BytesIO | None:
    if not url:
        return None
    try:
        resp = requests.get(
            url,
            timeout=30,
            headers={
                "User-Agent": _USER_AGENT,
                "Referer": "https://commons.wikimedia.org/",
                "Accept": "image/*,*/*;q=0.8",
            },
        )
        resp.raise_for_status()
    except requests.RequestException:
        return None

    # Guard against non-image payloads (e.g. an HTML rate-limit/error page).
    content_type = resp.headers.get("Content-Type", "").lower()
    if content_type and not content_type.startswith("image/"):
        return None
    if not resp.content:
        return None
    return io.BytesIO(resp.content)


def _get_placeholder(slide, name):
    for shape in slide.placeholders:
        if shape.name == name:
            return shape
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


@router.post("/docx")
def export_docx(req: WordExportRequest) -> StreamingResponse:
    require_session(req.session_id)
    base_name = _build_base_filename(req.citation)

    try:
        template_file = importlib.resources.files("highlight.data").joinpath(
            "highlight_template.docx"
        )
        template = DocxTemplate(template_file)

        photo_stream = _fetch_image_bytes(
            req.selected_image.full_url if req.selected_image else None
        )
        photo_link = req.selected_image.page_url if req.selected_image else ""

        context = {
            "title": req.title,
            "subtitle": req.subtitle,
            "photo": None,
            "photo_link": photo_link,
            "photo_site_name": _build_photo_site_name(req),
            "image_caption": req.image_caption,
            "science": req.science,
            "impact": req.impact,
            "summary": req.summary,
            "funding": req.funding,
            "citation": req.citation,
            "related_links": req.related_links,
            "point_of_contact": req.point_of_contact,
        }

        if photo_stream is not None and photo_stream.getbuffer().nbytes > 0:
            photo_stream.seek(0)
            context["photo"] = InlineImage(template, photo_stream, width=Mm(120))

        template.render(context)
        bio = io.BytesIO()
        template.save(bio)
        bio.seek(0)
    except Exception as exc:  # noqa: BLE001
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error generating Word document: {exc}",
        ) from exc

    return StreamingResponse(
        bio,
        media_type=_DOCX_MIME,
        headers={"Content-Disposition": f'attachment; filename="{base_name}.docx"'},
    )


@router.post("/pptx")
def export_pptx(req: PptExportRequest) -> StreamingResponse:
    session = require_session(req.session_id)
    base_name = _build_base_filename(req.citation)

    if not (req.title and req.objective and req.impact_points and req.approach_points):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=(
                "Please provide the title, objective, impact, and approach "
                "content before exporting."
            ),
        )

    figure_image = _resolve_pdf_image(session, req.figure_image_index)

    try:
        template_file = importlib.resources.files("highlight.data").joinpath(
            "highlight_template.pptx"
        )
        prs = Presentation(template_file)

        approach_points = req.approach_points or ["Approach not generated."]
        impact_points = req.impact_points or ["Impact points not generated."]

        for slide in prs.slides:
            impact_ph = _get_placeholder(slide, "Text Placeholder 10")
            approach_ph = _get_placeholder(slide, "Text Placeholder 9")
            picture_ph = _get_placeholder(slide, "Picture Placeholder 2")
            caption_ph = _get_placeholder(slide, "Text Placeholder 3")
            citation_ph = _get_placeholder(slide, "Text Placeholder 11")
            objective_ph = _get_placeholder(slide, "Text Placeholder 8")
            title_ph = _get_placeholder(slide, "Title 1")

            for ph in (
                title_ph,
                objective_ph,
                caption_ph,
                citation_ph,
                approach_ph,
                impact_ph,
                picture_ph,
            ):
                _remove_placeholder_outline(ph)

            if title_ph is not None:
                title_ph.text_frame.clear()
                title_ph.text_frame.text = req.title
            if objective_ph is not None:
                objective_ph.text_frame.clear()
                objective_ph.text_frame.text = req.objective
            if caption_ph is not None:
                caption_ph.text_frame.clear()
                caption_ph.text_frame.text = req.figure_caption
            if citation_ph is not None:
                citation_ph.text_frame.clear()
                citation_ph.text_frame.text = req.citation

            _populate_bullets(approach_ph, approach_points)
            _populate_bullets(impact_ph, impact_points)

            if picture_ph is not None and figure_image is not None:
                try:
                    figure_image.seek(0)
                    picture_ph.insert_picture(figure_image)
                except Exception:  # noqa: BLE001 - leave placeholder if insert fails
                    pass

        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
    except Exception as exc:  # noqa: BLE001
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"An error occurred while generating the PowerPoint: {exc}",
        ) from exc

    return StreamingResponse(
        ppt_io,
        media_type=_PPTX_MIME,
        headers={"Content-Disposition": f'attachment; filename="{base_name}.pptx"'},
    )


def _resolve_pdf_image(session: Session, index: int | None) -> io.BytesIO | None:
    if index is None or not session.pdf_bytes:
        return None
    extracted = hlt.extract_images_from_pdf(session.pdf_bytes)
    for item in extracted:
        if item["index"] == index:
            return io.BytesIO(item["bytes"])
    return None


def _remove_placeholder_outline(placeholder) -> None:
    """Force a placeholder's outline to 'no line'.

    PowerPoint renders a default dashed/dotted border around placeholders that
    do not explicitly declare a line style. Inserting an ``<a:ln><a:noFill/>``
    element on the shape properties suppresses that border.
    """
    if placeholder is None:
        return
    try:
        sp_pr = placeholder._element.spPr
    except AttributeError:
        return
    if sp_pr is None:
        return

    # Remove any existing line definition so we can set a clean 'no fill'.
    for existing in sp_pr.findall(qn("a:ln")):
        sp_pr.remove(existing)

    ln = sp_pr.makeelement(qn("a:ln"), {})
    ln.append(sp_pr.makeelement(qn("a:noFill"), {}))
    sp_pr.append(ln)


def _populate_bullets(placeholder, points) -> None:
    if placeholder is None or not hasattr(placeholder, "text_frame") or not points:
        return
    tf = placeholder.text_frame
    tf.clear()
    if len(tf.paragraphs):
        first = tf.paragraphs[0]
        if not first.text.strip() and len(first.runs) == 0:
            first._element.getparent().remove(first._element)
    for point_text in points[:3]:
        para = tf.add_paragraph()
        para.text = point_text.strip().lstrip("- ")
        para.level = 0
        para.font.size = Pt(13)
        para.alignment = PP_ALIGN.LEFT
