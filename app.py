import datetime
import io
import os
import importlib
import re 

from docxtpl import DocxTemplate, InlineImage
from docx.shared import Mm, Inches
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor 
import streamlit as st
# from openai import OpenAI
from langchain_openai import AzureChatOpenAI, OpenAI
import requests


import highlight as hlt
from highlight.utils import ApproachPoints, PydanticOutputParser, ImpactPoints
import highlight.prompts as prompts


def sanitize_filename(name):
    """Removes or replaces characters invalid for filenames."""
    # Remove characters not alphanumeric, underscore, hyphen, or period
    name = re.sub(r'[^\w\-\.]', '_', name)
    # Replace multiple consecutive underscores
    name = re.sub(r'_+', '_', name)
    # Remove leading/trailing underscores/periods
    name = name.strip('_.')
    # Limit length (optional but good practice)
    return name[:100] if len(name) > 100 else name


def get_placeholder(slide, name):
    for shape in slide.placeholders:
        # Check placeholder name (shape.name)
        if shape.name == name:
            return shape
    # Fallback: Check shapes by name if not found in placeholders (might be less reliable)
    for shape in slide.shapes:
            if shape.name == name:
                return shape
    return None


if "reduce_document" not in st.session_state:
    st.session_state.reduce_document = False

if "content_dict" not in st.session_state:
    st.session_state.content_dict = {}

if "model" not in st.session_state:
    st.session_state.model = "gpt-4o"

if "package" not in st.session_state:
    st.session_state.package = "langchain_azure_openai"

if "access" not in st.session_state:
    st.session_state.access = False 

# parameters for word document
if "title_response" not in st.session_state:
    st.session_state.title_response = None

if "subtitle_response" not in st.session_state:
    st.session_state.subtitle_response = None

if "photo" not in st.session_state:
    st.session_state.photo = None

if "photo_link" not in st.session_state:
    st.session_state.photo_link = None

if "photo_site_name" not in st.session_state:
    st.session_state.photo_site_name = None

if "image_caption" not in st.session_state:
    st.session_state.image_caption = None

if "science_response" not in st.session_state:
    st.session_state.science_response = None

if "impact_response" not in st.session_state:
    st.session_state.impact_response = None

if "summary_response" not in st.session_state:
    st.session_state.summary_response = None

if "funding" not in st.session_state:
    st.session_state.funding = None

if "citation" not in st.session_state:
    st.session_state.citation = None

if "base_export_filename" not in st.session_state:
    st.session_state.base_export_filename = f"ber-highlight_{datetime.date.today().strftime('%d%b%Y').lower()}"

if "related_links" not in st.session_state:
    st.session_state.related_links = None

# additional word doc content that is not in the template
if "figure_response" not in st.session_state:
    st.session_state.figure_response = None

if "figure_caption" not in st.session_state:
    st.session_state.figure_caption = None

if "figure_data" not in st.session_state:
    # Stores dict like {'Figure 1': 'Description...'} extracted from the paper for PPT selection
    st.session_state.figure_data = None

if "selected_figure_id" not in st.session_state:
    # Stores the ID (e.g., 'Figure 1') selected from the paper for the PPT slide
    st.session_state.selected_figure_id = None

if "caption_response" not in st.session_state:
    st.session_state.caption_response = None

if "output_file" not in st.session_state:
    st.session_state.output_file = None

# parameters for the ppt slide
if "objective_response" not in st.session_state:
    st.session_state.objective_response = None

if "approach_response" not in st.session_state:
    st.session_state.approach_response = None

if "ppt_impact_response" not in st.session_state:
    st.session_state.ppt_impact_response = None

if "figure_recommendation" not in st.session_state:
    st.session_state.figure_recommendation = None

if "search_phrase" not in st.session_state:
    st.session_state.search_phrase = None

if "point_of_contact" not in st.session_state:
    st.session_state.point_of_contact = None

if "project_info" not in st.session_state:
    st.session_state.project_info = {
        os.getenv("IM3_ACCESS", default=None): {
            "key": os.getenv("IM3_AZURE_OPENAI_API_KEY", default=None),
            "endpoint": os.getenv("IM3_AZURE_OPENAI_ENDPOINT", default=None),
            "deployment": "gpt-4o",
            "version": "2024-02-01",
            "project": "IM3"
        },
    }

if "active_project" not in st.session_state:
    st.session_state.active_project = "Other"

if "project_dict" not in st.session_state:
    st.session_state.project_dict = {
        "IM3": "Jennie Rice\nIM3 Principal Investigator\njennie.rice@pnnl.gov",
        "GCIMS": "Marshall Wise\nGCIMS Principal Investigator\nmarshall.wise@pnnl.gov",
        "COMPASS-GLM": "Robert Hetland\nCOMPASS-GLM Principal Investigator\nrobert.hetland@pnnl.gov",
        "ICoM": "Ian Kraucunas\nICoM Principal Investigator\nian.kraucunas@pnnl.gov",
        "Puget Sound": "Ning Sun\nPuget Sound Scoping and Pilot Study Principal Investigator\nning.sun@pnnl.gov",
        "Other": "First and Last Name\nCorresponding Project Name with POC Credentials\nEmail Address",
    }

# Figure selection state
if "figure_list" not in st.session_state:
    st.session_state.figure_list = None # Will hold the list ['Figure 1', 'Fig 2', ...]

if "selected_figure" not in st.session_state:
    st.session_state.selected_figure = None # Will hold the user's choice, e.g., 'Figure 1'

if "selected_figure_caption" not in st.session_state:
    st.session_state.selected_figure_caption = "" # Will hold the generated caption for the selected figure

if "wikimedia_query" not in st.session_state:
    st.session_state.wikimedia_query = ""

if "wikimedia_results" not in st.session_state:
    st.session_state.wikimedia_results = None # List of image dicts

if "selected_wikimedia_image_info" not in st.session_state:
    st.session_state.selected_wikimedia_image_info = None # Dict of selected image

if "wikimedia_limit" not in st.session_state:
    st.session_state.wikimedia_limit = 5

# States related to the chosen photo for the Word doc
if "photo" not in st.session_state: st.session_state.photo = None # Will hold BytesIO
if "photo_link" not in st.session_state: st.session_state.photo_link = None
if "photo_site_name" not in st.session_state: st.session_state.photo_site_name = None

if "extracted_pdf_images" not in st.session_state:
    # List of dicts like [{"index": 0, "page": 1, "xref": 123, "bytes": b'...'}, ...]
    st.session_state.extracted_pdf_images = None
if "ppt_figure_image_bytes" not in st.session_state:
    # Holds the raw bytes of the image chosen via "Assign Image" button
    st.session_state.ppt_figure_image_bytes = None

if "suggested_search_strings" not in st.session_state:
    st.session_state.suggested_search_strings = None

# ------------------------------------------------
# -- BEGIN INTERFACE --> 
# ------------------------------------------------

# Force responsive layout for columns also on mobile
st.write(
    """<style>
    [data-testid="column"] {
        width: calc(50% - 1rem);
        flex: 1 1 calc(50% - 1rem);
        min-width: calc(50% - 1rem);
    }
    </style>""",
    unsafe_allow_html=True,
)

st.markdown(
    """<h1 style='text-align: center;'>
        <span style='font-size:40px;'>&#128220;</span>  PAIGE  <span style='font-size:40px;'>&#128220;</span>
    </h1>
    <h3 style='text-align: center;'>The Pnnl AI assistant for GEnerating publication highlights</h3>
    <h5 style='text-align: center;'>Go from publication to a first draft highlight <i>fast</i>!</h5>
    """,
     unsafe_allow_html=True
)

with st.expander("**How to Use PAIGE**", expanded=False):
    st.markdown((
        "Simply: \n" + 
        "1. Enter in your project password or OpenAI API key \n"
        "2. Load the PDF document of your publication into the app \n" +  
        "3. Generate each part of your document in order \n" + 
        "4. Export the document to your local machine \n" + 
        "5. Repeat to generate the PowerPoint slide as well \n" + 
        "\n :memo: Note: Some parts of this process were left to be semi-automated. " + 
        "These include finding images that are free and open to use from a reliable " + 
        "source and choosing which figure from the paper to use in the PowerPoint slide. " + 
        "But don't worry, PAIGE offers helpers along the way."
    ))

if st.session_state.model in (["gpt-4o"]):
    st.session_state.max_allowable_tokens = 150000

# validate project and access key
if st.session_state.access is False:
    user_input = st.text_input(
        "Enter your project password or API key:", 
        type="password",
    )

    if user_input:
        if user_input in st.session_state.project_info.keys():
            project_info = st.session_state.project_info[user_input]
            st.session_state.active_project = project_info["project"]

            if st.session_state.package == "langchain_azure_openai":

                # setup environment locally
                os.environ["OPENAI_API_TYPE"]="azure"
                os.environ["OPENAI_API_VERSION"]=project_info["version"]
                os.environ["OPENAI_API_BASE"]=project_info["endpoint"]
                os.environ["OPENAI_API_KEY"]=project_info["key"]
                os.environ["OPENAI_CHAT_MODEL"]=project_info["deployment"]

                st.success(f"Hello {st.session_state.active_project} representative!", icon="✅")
                st.session_state.access = True

                st.session_state.client = AzureChatOpenAI(
                    deployment_name=project_info["deployment"],
                    azure_endpoint=project_info["endpoint"]
                )

        else:
            st.error(f"Invalid key or password.  Please provide a valid entry.", icon="🚨")
            st.session_state.access = False
            user_input = False

if st.session_state.access:

    st.markdown("### Upload file to process:")
    uploaded_file = st.file_uploader(
        label="### Select PDF or text file to upload",
        type=["pdf", "txt"],
        help="Select PDF or text file to upload",
    )

    if uploaded_file is not None:

        if uploaded_file.type == "text/plain":
            content_dict = hlt.read_text(uploaded_file)

        elif uploaded_file.type == "application/pdf":
            content_dict = hlt.read_pdf(uploaded_file)

        st.session_state.output_file = uploaded_file.name

        st.code(f"""File specs:\n
        - Number of pages:  {content_dict['n_pages']}
        - Number of characters:  {content_dict['n_characters']}
        - Number of words: {content_dict['n_words']}
        - Number of tokens: {content_dict['n_tokens']}
        """)

        if content_dict['n_tokens'] > st.session_state.max_allowable_tokens:
            msg = f"""
        The number of tokens in your document exceeds the maximum allowable tokens.
        This will cause your queries to fail.
        The queries account for the number of tokens in a prompt + the number of tokens in your document.
        
        Maximum allowable token count: {st.session_state.max_allowable_tokens}
        
        Your documents token count: {content_dict['n_tokens']}
        
        Token deficit: {content_dict['n_tokens'] - st.session_state.max_allowable_tokens}
        """
            st.error(msg, icon="🚨")

            st.session_state.reduce_document = st.radio(
                """Would you like me to attempt to reduce the size of 
            your document by keeping only relevant information? 
            If so, I will give you a file to download with the content 
            so you only have to do this once.
            If you choose to go through with this, it may take a while
            to process, usually on the order of 15 minutes for a 20K token
            document.
            Alternatively, you can copy and paste the contents that you
            know are of interest into a text file and upload that
            instead.
        
            """,
                ("Yes", "No"),
            )

        # word document content
        st.markdown("### Section 1:  Content to fill in Word document template:")

# ------------------------------------------------
# -- DOC:  START TITLE SECTION --> 
# ------------------------------------------------

        st.markdown("---")

        # title section
        title_container = st.container()
        title_container.markdown("##### Generate title from text content")

        # title criteria
        title_container.markdown("""
        The title should meet the following criteria:
        - No colons are allowed in the output.
        - Should pique the interest of the reader while still being somewhat descriptive.
        - Be understandable to a general audience.
        - Should be only once sentence.
        - Should have a maximum length of 10 words.
        """)

        title_container.markdown("Set desired temperature:")

        # title slider
        title_temperature = title_container.slider(
            "Title Temperature",
            0.0,
            1.0,
            0.2,
            label_visibility="collapsed"
        )

        # build container content
        if title_container.button('Generate Title'):

            st.session_state.title_response = hlt.generate_content(
                client=st.session_state.client,
                container=title_container,
                content=content_dict["content"],
                prompt_name="title",
                result_title="Title Result:",
                max_tokens=50,
                temperature=title_temperature,
                box_height=75,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model
            )

        else:
            if st.session_state.title_response is not None:
                title_container.markdown("Title Result:")
                title_container.text_area(
                    label="Title Result:",
                    value=st.session_state.title_response,
                    label_visibility="collapsed",
                    height=75
                )

# ------------------------------------------------
# -- DOC:  START SUBTITLE SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # subtitle section
        subtitle_container = st.container()
        subtitle_container.markdown("##### Generate subtitle from text content")

        # subtitle criteria
        subtitle_container.markdown("""
        The subtitle should meet the following criteria:
        - Be an extension of and related to, but not directly quote, the title.
        - Provide information that will make the audience want to find out more about the research.
        - Do not use more than 155 characters including spaces.
        """)

        subtitle_container.markdown("Set desired temperature:")

        # subtitle slider
        subtitle_temperature = subtitle_container.slider(
            "Subtitle Temperature",
            0.0,
            1.0,
            0.5,
            label_visibility="collapsed"
        )

        # build container content
        if subtitle_container.button('Generate Subtitle'):

            if st.session_state.title_response is None:
                st.write("Please generate a Title first.  Subtitle generation considers the title response.")
            else:

                st.session_state.subtitle_response = hlt.generate_content(
                    client=st.session_state.client,
                    container=subtitle_container,
                    content=content_dict["content"],
                    prompt_name="subtitle",
                    result_title="Subtitle Result:",
                    max_tokens=100,
                    temperature=subtitle_temperature,
                    box_height=75,
                    additional_content=st.session_state.title_response,
                    max_word_count=100,
                    min_word_count=75,
                    max_allowable_tokens=st.session_state.max_allowable_tokens,
                    model=st.session_state.model
                )

        else:
            if st.session_state.subtitle_response is not None:
                subtitle_container.markdown("Subtitle Result:")
                subtitle_container.text_area(
                    label="Subtitle Result:",
                    value=st.session_state.subtitle_response,
                    label_visibility="collapsed",
                    height=75
                )

# ------------------------------------------------
# -- DOC:  START SCIENCE SUMMARY SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # science section
        science_container = st.container()
        science_container.markdown("##### Generate science summary from text content")

        # science criteria
        science_container.markdown("""
        **GOAL**:  Describe the scientific results for a non-expert, non-scientist audience.
        
        The description should meet the following criteria:
        - Answer what the big challenge in this field of science is that the research addresses.
        - State what the key finding is.
        - Explain the science, not the process.
        - Be understandable to a high school senior or college freshman.
        - Use short sentences and succinct words.
        - Avoid technical terms if possible.  If technical terms are necessary, define them.
        - Provide the necessary context so someone can have a very basic understanding of what you did. 
        - Start with topics that the reader already may know and move on to more complex ideas.
        - Use present tense.
        - In general, the description should speak about the research or researchers in first person.
        - Use a minimum of 75 words and a maximum of 100 words. 
        """)

        science_container.markdown("Set desired temperature:")

        # slider
        science_temperature = science_container.slider(
            "Science Summary Temperature",
            0.0,
            1.0,
            0.3,
            label_visibility="collapsed"
        )

        # build container content
        if science_container.button('Generate Science Summary'):
            st.session_state.science_response = hlt.generate_content(
                client=st.session_state.client,
                container=science_container,
                content=content_dict["content"],
                prompt_name="science",
                result_title="Science Summary Result:",
                max_tokens=200,
                temperature=science_temperature,
                box_height=250,
                max_word_count=100,
                min_word_count=75,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model
            )

        else:
            if st.session_state.science_response is not None:
                science_container.markdown("Science Summary Result:")
                science_container.text_area(
                    label="Science Summary Result:",
                    value=st.session_state.science_response,
                    label_visibility="collapsed",
                    height=250
                )

# ------------------------------------------------
# -- DOC:  START IMPACT SUMMARY SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # impact section
        impact_container = st.container()
        impact_container.markdown("##### Generate impact summary from text content")

        impact_container.markdown("""
        **GOAL**: Describe the impact of the research to a non-expert, non-scientist audience.
        
        The description should meet the following criteria:
        - Answer why the findings presented are important, i.e., what problem the research is trying to solve.
        - Answer if the finding is the first of its kind.
        - Answer what was innovative or distinct about the research.
        - Answer what the research enables other scientists in your field to do next.
        - Include other scientific fields potentially impacted. 
        - Be understandable to a high school senior or college freshman. 
        - Use short sentences and succinct words.
        - Avoid technical terms if possible.  If technical terms are necessary, define them.
        - Use present tense.
        - In general, the description should speak about the research or researchers in first person.
        - Use a minimum of 75 words and a maximum of 100 words. 
        """)


        impact_container.markdown("Set desired temperature:")

        # slider
        impact_temperature = impact_container.slider(
            "Impact Summary Temperature",
            0.0,
            1.0,
            0.0,
            label_visibility="collapsed"
        )

        # build container content
        if impact_container.button('Generate Impact Summary'):
            st.session_state.impact_response = hlt.generate_content(
                client=st.session_state.client,
                container=impact_container,
                content=content_dict["content"],
                prompt_name="impact",
                result_title="Impact Summary Result:",
                max_tokens=700,
                temperature=impact_temperature,
                box_height=250,
                max_word_count=100,
                min_word_count=75,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model
            )

        else:
            if st.session_state.impact_response is not None:
                impact_container.markdown("Impact Summary Result:")
                impact_container.text_area(
                    label="Impact Summary Result:",
                    value=st.session_state.impact_response,
                    label_visibility="collapsed",
                    height=250
                )

# ------------------------------------------------
# -- DOC:  START GENERAL SUMMARY SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # general summary section
        summary_container = st.container()
        summary_container.markdown("##### Generate general summary from text content")

        summary_container.markdown("""
        **GOAL**: Generate a general summary of the current research.
        
        The summary should meet the following criteria:
        - Should relay key findings and value.
        - The summary should be still accessible to the non-specialist but may be more technical if necessary. 
        - Do not mention the names of institutions. 
        - If there is a United States Department of Energy Office of Science user facility involved, such as NERSC, you can mention the user facility. 
        - Should be 1 or 2 paragraphs detailing the research.
        - Use present tense.
        - In general, the description should speak about the research or researchers in first person.
        - Use no more than 200 words.
        """)

        summary_container.markdown("Set desired temperature:")

        # slider
        summary_temperature = summary_container.slider(
            "General Summary Temperature",
            0.0,
            1.0,
            0.3,
            label_visibility="collapsed"
        )

        # build container content
        if summary_container.button('Generate General Summary'):
            st.session_state.summary_response = hlt.generate_content(
                client=st.session_state.client,
                container=summary_container,
                content=content_dict["content"],
                prompt_name="summary",
                result_title="General Summary Result:",
                max_tokens=700,
                temperature=summary_temperature,
                box_height=400,
                max_word_count=200,
                min_word_count=100,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model
            )

        else:
            if st.session_state.summary_response is not None:
                summary_container.markdown("General Summary Result:")
                summary_container.text_area(
                    label="General Summary Result:",
                    value=st.session_state.summary_response,
                    label_visibility="collapsed",
                    height=400
                )

# ------------------------------------------------
# -- DOC:  START CITATION SELECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # citation recommendations section
        citation_container = st.container()
        citation_container.markdown("##### Citation for the paper in Chicago style")
        citation_container.markdown("This will only use what is represented in the publication provided.")

        if citation_container.button('Generate Citation'):
            st.session_state.citation = hlt.generate_content(
                client=st.session_state.client,
                container=citation_container,
                content=content_dict["content"],
                prompt_name="citation",
                result_title="",
                max_tokens=300,
                temperature=0.0,
                box_height=200,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model
            ).replace('"', "")

        else:
            if st.session_state.citation is not None:
                citation_container.text_area(
                    label="Citation",
                    value=st.session_state.citation,
                    label_visibility="collapsed",
                    height=200
                )

# ------------------------------------------------
# -- DOC:  START PHOTO SELECTION --> 
# ------------------------------------------------

        st.markdown("---") 

        st.markdown("##### Find an Image for your Word Document")
        st.markdown("**Note**:  This is not an image from your paper.  It is meant to be an editorial cover image.")
        st.markdown(
            "The following is a convenience service and uses Wikimedia Commons due to their images being fully open and reusable.  " + 
            "That being stated, you may not always be able to find what you need in their database and may have to " + 
            "search another resource.  Simply do not execute this block if you intend to find an image elsewhere."
        )
        # --- Main Container for this Section ---
        img_search_container = st.container(border=True)

        # Step 1 - Suggest Search Strings
        img_search_container.markdown("##### 1. Generate Suggested Wikimedia Search Strings (Optional)")
        img_search_container.caption(
            "These are produced from the 'General Summary' generated earlier.  " + 
            "Simple searches are more productive, so use these as general guidance.  " + 
            "\n\n**For example**, simply searching for 'Groundwater irrigation' will work better than " + 
            "'Groundwater irrigation in the Snake River Basin'."
        )

        suggest_container = img_search_container.container() # Sub-container for this step

        # Check if summary exists first
        if st.session_state.summary_response:
            if suggest_container.button("Suggest Search Strings", key="suggest_wiki_search"):
                with st.spinner("Generating search string ideas..."):
                    # Call generate_content for suggestions
                    # Assuming 'figure' prompt takes summary and returns newline-separated strings
                    suggestions = hlt.generate_content(
                        client=st.session_state.client,
                        container=suggest_container, # Display result within this container
                        content=st.session_state.summary_response, # Input is the summary
                        prompt_name="figure", # Use the prompt for generating search strings
                        result_title="Suggested Strings:",
                        max_tokens=200, # Adjust as needed
                        temperature=0.5, # Moderate temperature for suggestions
                        box_height=150, # Text area height
                        # Assuming generate_content displays result in text_area, no max/min words needed here
                        max_allowable_tokens=st.session_state.max_allowable_tokens,
                        model=st.session_state.model,
                        package=st.session_state.package
                    )
                    # Store the raw response (might be newline separated string)
                    st.session_state.suggested_search_strings = suggestions
                    # Rerun not strictly necessary as generate_content handles display,
                    # but needed if default query logic below should immediately update
                    st.rerun()

            # Display existing suggestions if already generated
            elif st.session_state.suggested_search_strings:
                suggest_container.markdown("**Suggested Strings:**")
                suggest_container.text_area(
                    label="Suggested Strings", # Hidden label
                    value=st.session_state.suggested_search_strings.replace('"', ''),
                    height=150,
                    disabled=False,
                    label_visibility="collapsed"
                )

        else:
            suggest_container.warning("Please generate the 'General Summary' in Section 1 first to enable suggestions.")


        img_search_container.markdown("---") # Separator

        # --- Step 2: Search and Select Image (Previously Step 1) ---
        img_search_container.markdown("##### 2. Search Wikimedia Commons and Select Image")

        # --- Determine Default Query ---
        default_query = ""
        # PRIORITIZE first suggested string if available
        if st.session_state.suggested_search_strings:
            first_suggestion = st.session_state.suggested_search_strings.split('\n')[0].strip()
            if first_suggestion:
                default_query = first_suggestion
        # Fallback to title or summary
        elif st.session_state.title_response:
            default_query = st.session_state.title_response
        elif st.session_state.summary_response:
            default_query = " ".join(st.session_state.summary_response.split()[:15])

        # --- Search Controls Row (No change needed here) ---
        controls_cols = img_search_container.columns([3, 1, 1])
        with controls_cols[0]: # Search Query Input
            user_query = st.text_input(
                "Image Search Query:",
                value=st.session_state.get("wikimedia_query", default_query), # Use default logic
                key="wikimedia_query_input"
            )
            st.session_state.wikimedia_query = user_query

        with controls_cols[1]: # Number of Images Input
            num_images = st.number_input(
                "Max Results:", min_value=3, max_value=30,
                value=st.session_state.wikimedia_limit, step=3,
                key="wikimedia_limit_input", help="Number of images to retrieve (max 30)"
            )
            st.session_state.wikimedia_limit = int(num_images)

        with controls_cols[2]: # Search Button
            st.write("") # Placeholder for spacing
            st.write("")
            search_button = st.button("Search", key="wiki_search_btn", use_container_width=True)


        # --- Clear Button ---
        if st.session_state.wikimedia_results or st.session_state.selected_wikimedia_image_info:
            if img_search_container.button("Clear Search / Reset Selection", key="wiki_clear_btn"):
                st.session_state.wikimedia_results = None
                st.session_state.selected_wikimedia_image_info = None
                st.session_state.photo = None
                st.session_state.photo_link = None
                st.session_state.photo_site_name = None
                st.rerun()

        # --- Execute Search ---
        if search_button and st.session_state.wikimedia_query:
            with st.spinner("Searching Wikimedia Commons..."):
                st.session_state.wikimedia_results = hlt.search_wikimedia_commons(
                    query=st.session_state.wikimedia_query,
                    limit=st.session_state.wikimedia_limit
                )
                st.session_state.selected_wikimedia_image_info = None # Reset previous selection
                if not st.session_state.wikimedia_results:
                    img_search_container.info("No suitable images found for your query.")
                # No rerun here, let results display immediately below
                # st.rerun() # Removing rerun for smoother feel

        # --- Display Results in Columns ---
        if st.session_state.wikimedia_results:

            if st.session_state.selected_wikimedia_image_info is None:

                num_found = len(st.session_state.wikimedia_results)
                # Title ABOVE the results box
                img_search_container.markdown(f"**Found {num_found} image(s):** (Select one below)")

                # --- Create a NEW sub-container specifically for the results grid ---
                results_grid_container = st.container(border=True, height=500)

                results_grid_container.write("##### 3. Select an image")

                # --- Place columns and results INSIDE the new sub-container ---
                with results_grid_container: # Use 'with' block for clarity
                    num_columns = 3 # Adjust as desired
                    cols = st.columns(num_columns) # Create columns inside the sub-container

                    for i, img_data in enumerate(st.session_state.wikimedia_results):
                        col_index = i % num_columns
                        with cols[col_index]: # Place content in the current column
                            if img_data.get("thumbnail_url"):
                                st.image(
                                    img_data["thumbnail_url"],
                                    caption=f"{img_data.get('title', 'N/A')} ({img_data.get('license', 'N/A')})",
                                    width=180 # Adjust width
                                )
                                if img_data.get('page_url'):
                                    st.caption(f"[View on Wikimedia]({img_data['page_url']})")

                                # Selection Button for each image
                                button_key = f"select_wiki_{img_data.get('id', i)}"
                                if st.button(f"Select This Image", key=button_key):
                                    # --- Selection logic (remains the same) ---
                                    st.session_state.selected_wikimedia_image_info = img_data
                                    with st.spinner("Preparing selected image..."):
                                        try:
                                            headers = {'User-Agent': 'PAIGE/1.0 (Highlight Generator App)'}
                                            img_response = requests.get(img_data['full_url'], stream=True, timeout=15, headers=headers)
                                            img_response.raise_for_status()
                                            img_bytes_io = io.BytesIO(img_response.content)
                                            st.session_state.photo = img_bytes_io
                                            st.session_state.photo_link = img_data.get('page_url', '')
                                            st.session_state.photo_site_name = "Wikimedia Commons"
                                            st.success(f"Image '{img_data.get('title', 'Selected')}' prepared for Word doc.")
                                            st.rerun()
                                        except Exception as e:
                                            st.error(f"Failed to download/prepare selected image: {e}")
                                            # Reset relevant states on failure
                                            st.session_state.photo = None
                                            st.session_state.photo_link = None
                                            st.session_state.photo_site_name = None
                                            st.session_state.selected_wikimedia_image_info = None
                                    # --- End selection logic ---
                                # Add separator below button
                                st.markdown("---")

                            else:
                                st.caption(f"Thumbnail missing for '{img_data.get('title', 'N/A')}'")
                                st.markdown("---")


        # --- Display Final Selection Info and Download Button ---
        if st.session_state.selected_wikimedia_image_info:
            # This section remains unchanged from the previous version
            img_search_container.markdown("---")
            img_search_container.markdown("##### 3. Selected Image:")
            selected_info = st.session_state.selected_wikimedia_image_info

            # Display image thumbnail and details
            sel_col1, sel_col2 = img_search_container.columns([1, 2])
            with sel_col1:
                if selected_info.get("thumbnail_url"):
                    st.image(selected_info["thumbnail_url"], width=150)
            with sel_col2:
                st.markdown(f"**Title:** {selected_info.get('title', 'N/A')}")
                st.markdown(f"**License:** {selected_info.get('license', 'N/A')}")
                if selected_info.get('page_url'):
                    st.markdown(f"**Attribution/Source:** [{selected_info['page_url']}]({selected_info['page_url']})")
                if selected_info.get('artist'):
                    st.caption(f"Artist Info: {selected_info['artist']}", unsafe_allow_html=True)
                st.info("Selected image & attribution link will be used in Word doc export.")

            # CAPTION GENERATION/EDIT SECTION
            img_search_container.markdown("---") # Separator
            img_search_container.markdown("##### 4. Generate a General Caption (based on paper summary)")
            # Using a sub-container for layout clarity, optional
            caption_editor_container = img_search_container.container()

            # Define parameters for caption generation
            caption_prompt_name = "figure_caption" # Assumes this prompt summarizes the paper
            caption_temperature = 0.3              # Adjust temperature if needed
            caption_max_tokens = 100               # Max tokens for the LLM response
            caption_box_height = 100               # Height of the text area
            caption_max_words = 30                 # Target max words (adjust if needed)
            caption_min_words = 10                 # Target min words

            # Button to generate a caption suggestion
            if caption_editor_container.button("Suggest Caption", key="gen_wiki_caption"):
                with st.spinner("Generating caption suggestion..."):
                    # We use generate_content which handles UI updates within the container
                    # Base the caption on the *paper's* content, not the image metadata
                    generated_caption = hlt.generate_content(
                        client=st.session_state.client,
                        container=caption_editor_container, # Place result inside this container
                        content=content_dict["content"],    # Use main paper content as input
                        prompt_name=caption_prompt_name,
                        result_title="Suggested Caption (Editable):", # Title for text area
                        max_tokens=caption_max_tokens,
                        temperature=caption_temperature,
                        box_height=caption_box_height,
                        max_word_count=caption_max_words,
                        min_word_count=caption_min_words,
                        max_allowable_tokens=st.session_state.max_allowable_tokens,
                        model=st.session_state.model,
                        package=st.session_state.package
                    )
                    # Store the generated caption in the correct state variable
                    st.session_state.image_caption = generated_caption
                    # Rerun needed to ensure the text_area below picks up the new value if generate_content doesn't update it seamlessly
                    st.rerun()

            # Display the current caption (if it exists) in an editable text area
            # This allows viewing generated caption or editing it, or entering one manually
            current_caption_value = st.session_state.image_caption if st.session_state.image_caption is not None else ""
            edited_caption = caption_editor_container.text_area(
                label="Image Caption (Editable):", # Label for screen readers etc.
                value=current_caption_value,
                key="edit_image_caption",
                height=caption_box_height,
                label_visibility="collapsed", # Hide label as markdown title is present
                placeholder="Enter caption or generate suggestion..."
            )

            # Update the session state if the user modifies the text area
            if edited_caption != current_caption_value:
                st.session_state.image_caption = edited_caption
                st.rerun() # Rerun to confirm the change persists visually


            # Download Button Logic
            img_search_container.markdown("---")
            img_search_container.markdown("##### 5. Download Full Resolution Image")
            mime_type = selected_info.get('mime', 'application/octet-stream')
            extension_map = {
                'image/jpeg': '.jpg', 'image/png': '.png', 'image/gif': '.gif',
                'image/svg+xml': '.svg', 'image/tiff': '.tif'}
            file_extension = extension_map.get(mime_type, '.png')
            base_filename = sanitize_filename(selected_info.get('title', 'wikimedia_image'))
            download_filename = f"{st.session_state.base_export_filename}{file_extension}"

            try:
                # --- Define headers for download ---
                headers = {'User-Agent': 'PAIGE/1.0 (Highlight Generator App Download)'} # Be descriptive
                # --- Fetch the image bytes for the download button ---
                with st.spinner(f"Preparing '{download_filename}'..."):
                    image_bytes = requests.get(
                        selected_info['full_url'],
                        timeout=30,
                        headers=headers # <-- Add headers argument
                    ).content

                # Display the download button
                img_search_container.download_button(
                label=f"Download {file_extension.upper()[1:]} Image",
                data=image_bytes,
                file_name=download_filename,
                mime=mime_type,
                key=f"download_{selected_info.get('id', 'selected')}"
                )
                img_search_container.caption(f"Selected file: `{download_filename}` ({mime_type})")
                img_search_container.caption(f"Remember to check attribution requirements at the source link above.")

            except requests.exceptions.RequestException as req_e:
                img_search_container.error(f"Error downloading image: {req_e}")
            except Exception as e:
                img_search_container.error(f"Could not prepare image for download: {e}")

# ------------------------------------------------
# -- DOC:  START FUNDING SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # funding recommendations section
        funding_container = st.container()
        funding_container.markdown("##### Funding statement from the paper")
        funding_container.markdown((
            "Note:  Some journals house this information in the sidebar of the PDF " + 
            "and add in another similar, but very different, statment in the text. " + 
            "Please ensure you review this statement thouroughly to prevent any errors."
        ))

        if funding_container.button('Generate Funding Statement'):
            st.session_state.funding = hlt.generate_content(
                client=st.session_state.client,
                container=funding_container,
                content=content_dict["content"],
                prompt_name="funding",
                result_title="",
                max_tokens=300,
                temperature=0.0,
                box_height=200,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model
            ).replace('"', "")

        else:
            if st.session_state.funding is not None:
                funding_container.text_area(
                    label="Funding statement",
                    value=st.session_state.funding,
                    label_visibility="collapsed",
                    height=200
                )

# ------------------------------------------------
# -- DOC:  START POINT OF CONTACT SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # point of contact box
        poc_container = st.container()
        poc_container.markdown("##### Point of contact for the research by project")

        # select the POC information from the dropdown
        st.session_state.point_of_contact = st.session_state.project_dict[
            poc_container.selectbox(
            label="Select the project who funded the work:",
            options=[
                st.session_state.active_project,
                "COMPASS-GLM", 
                "GCIMS", 
                "ICoM",
                "IM3",
                "Puget Sound",
                "Other",
            ])
        ]

        
        poc_container.write("What will be written to the document as the point of contact:")
        poc_parts = st.session_state.point_of_contact.split("\n")
        poc_container.success(
            f"""
            {poc_parts[0]}\n
            {poc_parts[1]}\n
            {poc_parts[2]}\n
            """
        )

# ------------------------------------------------
# -- Generate base file name for exports --> 
# ------------------------------------------------

        if st.session_state.citation:
            citation = st.session_state.citation
            # Get today's date formatted
            today_str = datetime.date.today().strftime("%d%b%Y").lower()

            # --- Attempt to parse citation ---
            # NOTE: This parsing is based on common formats and might fail for others.
            last_name = "unknown"
            year = "YYYY"
            journal_abbrev = "journal"

            try:
                # 1. Extract First Author Last Name
                # Assumes format "LastName, F. M., Second Author..." or "LastName FM, Second Author..."
                first_author_part = citation.split(',')[0].strip()
                # Handle cases like "LastName FM" vs "LastName"
                if ' ' in first_author_part and not first_author_part.isupper(): # Avoid splitting acronyms/all caps
                    last_name = first_author_part.split(' ')[0].strip().lower()
                else:
                    last_name = first_author_part.lower()
                # Basic cleaning if needed
                last_name = re.sub(r'[^a-z\-]', '', last_name) # Keep letters, hyphen
                if not last_name: last_name = "unknown" # Fallback

            except Exception:
                last_name = "unknown" # Fallback on any parsing error

            try:
                # 2. Extract Year (look for 4 digits often after author block or date)
                # Regex looks for ". YYYY." or "(YYYY)." or " YYYY;" etc.
                match = re.search(r'[\.\(;]\s*(\d{4})[\.\);]', citation)
                if match:
                    year = match.group(1)
                else: # Fallback if first pattern fails
                    match = re.search(r'\b(\d{4})\b', citation) # Find any 4-digit number
                    if match: year = match.group(1) # Less reliable, might pick wrong number

            except Exception:
                year = "YYYY"

            try:
                # 3. Extract Journal Abbreviation (This is the most fragile part)
                # Attempt 1: Look for text between title (in quotes) and volume (digits)
                match = re.search(r'["”]\s*\.\s*([^,.:]+?)\s*\d+[:\(]', citation, re.IGNORECASE)
                if match:
                    journal_part = match.group(1).strip()
                    # Abbreviate: lowercase, remove spaces/periods/commas
                    journal_abbrev = re.sub(r'[\s\.,]', '', journal_part).lower()
                else:
                    # Attempt 2: Look for likely candidates based on common knowledge (less robust)
                    # This part could be expanded significantly or use external lookups
                    if "applied energy" in citation.lower(): journal_abbrev = "appliedenergy"
                    elif "journal of" in citation.lower(): journal_abbrev = "journal" # Example placeholder
                    # Add more rules or default
                    else: journal_abbrev = "journal" # Default if not found

                if not journal_abbrev: journal_abbrev = "journal" # Final fallback

            except Exception:
                journal_abbrev = "journal"

            # --- Construct the base filename ---
            st.session_state.base_export_filename = f"{last_name}_etal_{year}_{journal_abbrev}_ber-highlight_{today_str}"

# ------------------------------------------------
# -- DOC:  START EXPORT SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        export_container = st.container()
        export_container.markdown("##### Export Word document with new content when ready")

        # template parameters
        word_parameters = {
            'title': st.session_state.title_response,
            'subtitle': st.session_state.subtitle_response,
            'photo': None,
            'photo_link': st.session_state.photo_link,
            'photo_site_name': "Wikimedia Commons",
            'image_caption': st.session_state.image_caption,
            'science': st.session_state.science_response,
            'impact': st.session_state.impact_response,
            'summary': st.session_state.summary_response,
            'funding': st.session_state.funding,
            'citation': st.session_state.citation,
            'related_links': st.session_state.related_links,
            'point_of_contact': st.session_state.point_of_contact,
        }

        # --- Construct Attribution String and Ensure Correct Link ---
        if st.session_state.selected_wikimedia_image_info:
            selected_info = st.session_state.selected_wikimedia_image_info

            # 1. Construct the full attribution string for 'photo_site_name'
            artist = selected_info.get('artist_plain', '')
            license_short = selected_info.get('license', '')
            license_url = selected_info.get('license_url', None)

            parts = []
            if artist and artist != "Unknown Artist":
                parts.append(artist)
            if license_short:
                license_part = license_short
                # Only add URL part if URL exists
                if license_url:
                    license_part += f" <{license_url}>"
                parts.append(license_part)
            parts.append("via Wikimedia Commons")

            # Assign the constructed string to 'photo_site_name' key
            word_parameters['photo_site_name'] = ", ".join(filter(None, parts)).replace("<", "(").replace(">", ")")

            # 2. Ensure 'photo_link' has the Wikimedia file page URL (it should already from selection logic)
            word_parameters['photo_link'] = st.session_state.photo_link # Or selected_info.get('page_url', '')


        # --- Load template ---
        try:
            word_template_file = importlib.resources.files('highlight.data').joinpath('highlight_template.docx')
            template = DocxTemplate(word_template_file)

            # --- Process Image for 'photo' placeholder (using InlineImage) ---
            if isinstance(st.session_state.photo, io.BytesIO) and st.session_state.photo.getbuffer().nbytes > 0:
                st.session_state.photo.seek(0)
                image_for_template = InlineImage(template, st.session_state.photo, width=Mm(120)) # Adjust width
                word_parameters['photo'] = image_for_template
            else:
                word_parameters['photo'] = None

            # --- Render the template ---
            template.render(word_parameters)
            bio = io.BytesIO()
            template.save(bio)

            # --- Provide Download Button ---
            export_container.download_button(
                label="Export Word Document",
                data=bio.getvalue(),
                file_name=f"{st.session_state.base_export_filename}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

        except Exception as e:
            export_container.error(f"Error generating Word document: {e}")
            export_container.error("Check template placeholders ({{photo}}, {{photo_link}}, {{photo_site_name}}) and image preparation.")

# ------------------------------------------------
# -- PPT:  START POWER POINT SECTION --> 
# ------------------------------------------------

        st.markdown("---")  

        # power point slide content
        st.markdown("### Section 2:  Content to fill in PowerPoint template:")

# ------------------------------------------------
# -- PPT:  START OBJECTIVE SECTION --> 
# ------------------------------------------------

        st.markdown("---") 

        # objective section
        objective_container = st.container()
        objective_container.markdown("##### Generate objective summary from text content")

        objective_container.markdown("""
        **GOAL**:  Generate one sentence stating the core purpose of the study.
        
        The sentence should meet the following criteria:
        - Use active verbs for the start of each point.
        - Use present tense.
        - Do not include methodology related to statistical, technological, and theory based
        """)

        objective_container.markdown("Set desired temperature:")

        # slider
        objective_temperature = objective_container.slider(
            "Objective Temperature",
            0.0,
            1.0,
            0.3,
            label_visibility="collapsed"
        )

        # build container content
        if objective_container.button('Generate Objective'):
            st.session_state.objective_response = hlt.generate_content(
                client=st.session_state.client,
                container=objective_container,
                content=content_dict["content"],
                prompt_name="objective",
                result_title="Objective Result:",
                max_tokens=300,
                temperature=objective_temperature,
                box_height=250,
                max_allowable_tokens=st.session_state.max_allowable_tokens,
                model=st.session_state.model
            )

        else:
            if st.session_state.objective_response is not None:
                objective_container.markdown("Objective Result:")
                objective_container.text_area(
                    label="Objective Result:",
                    value=st.session_state.objective_response,
                    label_visibility="collapsed",
                    height=250
                )

# ------------------------------------------------
# -- PPT:  START APPROACH SECTION --> 
# ------------------------------------------------

        st.markdown("---") 

        approach_container = st.container()
        approach_container.markdown("##### Generate approach summary from text content")

        approach_container.markdown("""
        **GOAL**:  Clearly and concisely state in 2-3 short points how this work accomplished the stated objective from a methodolgocial perspecive.
        - Based off of the objective summary 
        - Only include methodology including but not limited to: statistical, technological, and theory based approaches. 
        - Use a different action verb to start sentences than what is used to begin the objective statement.
        - Use active verbs for the start of each point.  
        - Use present tense.
        """)

        approach_container.markdown("Set desired temperature:")

        # slider
        approach_temperature = approach_container.slider(
            "Approach Temperature",
            0.0,
            1.0,
            0.1,
            label_visibility="collapsed"
        )

        # build container content
        if approach_container.button('Generate Approach'):
            if st.session_state.objective_response is None:
                approach_container.warning("Please generate the Objective first.")
            else:
                with st.spinner("Generating approach points..."):
                    try:
                        # 1. Get the user prompt string
                        user_prompt_for_approach = hlt.generate_prompt(
                            content=content_dict["content"],
                            prompt_name="approach",
                            additional_content=st.session_state.objective_response
                        )

                        # 2. Instantiate the parser
                        parser = PydanticOutputParser(pydantic_object=ApproachPoints)

                        # 3. Call the structured generation function
                        structured_result = hlt.generate_structured_content(
                            client=st.session_state.client,
                            system_scope=prompts.SYSTEM_SCOPE,
                            user_prompt=user_prompt_for_approach,
                            pydantic_parser=parser,
                            max_tokens=300, # Adjust as needed
                            temperature=approach_temperature,
                            max_allowable_tokens=st.session_state.max_allowable_tokens,
                            model=st.session_state.model,
                            package=st.session_state.package
                        )

                        # 4. Store the list of points
                        st.session_state.approach_response = structured_result.points
                        # approach_container.success("Approach points generated!")

                    except Exception as e:
                        st.session_state.approach_response = None
                        # approach_container.error(f"Failed to generate approach: {e}")

         # Display the approach points in an editable text box
        if st.session_state.approach_response is not None: # Check if it exists (could be an empty list)
            approach_container.markdown("Approach Result:")
            response_data = st.session_state.approach_response # Should be a list of strings

            # --- Prepare string value for the text_area ---
            if isinstance(response_data, list):
                # Format the list into a multi-line string with bullets
                # Ensure points don't already start with '-' before adding one
                # Include only non-empty points
                current_text_value = "\n".join([
                    f"{str(point).strip()}" if str(point).strip().startswith("-") else f"- {str(point).strip()}"
                    for point in response_data if str(point).strip()
                ])
            elif isinstance(response_data, str):
                # Fallback if it somehow received a string (shouldn't happen ideally)
                current_text_value = response_data
                # approach_container.warning("Approach data was a string, expected a list. Displaying as is.")
            else:
                # Handle unexpected types
                approach_container.error(f"Cannot display approach result: Unexpected data type {type(response_data)}")
                current_text_value = "" # Default to empty

            # --- Use st.text_area for display and editing ---
            edited_text = approach_container.text_area(
                label="Approach Result (Editable):", # Label for accessibility
                value=current_text_value,
                height=250, # Adjust height as needed
                key="approach_edit_area", # Add a unique key
                label_visibility="collapsed" # Hide label visually
            )

            # --- Update session state with the potentially edited content ---
            # Parse the edited text FROM the text area back into a LIST
            updated_approach_list = [
                line.strip().lstrip('- ') # Remove leading bullet and space
                for line in edited_text.split('\n') if line.strip() # Split lines, ignore empty
            ]

            # Store the updated list back into session state
            # This ensures consistency for PowerPoint generation etc.
            st.session_state.approach_response = updated_approach_list

# ------------------------------------------------
# -- PPT:  START IMPACT SECTION --> 
# ------------------------------------------------

        st.markdown("---")

        # power point impact section
        ppt_impact_container = st.container()
        ppt_impact_container.markdown("##### Generate impact points from text content")

        ppt_impact_container.markdown("""
        **GOAL**:  Clearly and concisely state in 3 points the key results and outcomes from this research. 
        - State what the results indicate.
        - Include results that may be considered profound or surprising.
        - Each point should be 1 concise sentence.
        - Use present tense.
        """
        )

        ppt_impact_container.markdown("Set desired temperature:")

        # slider
        ppt_impact_temperature = ppt_impact_container.slider(
            "Impact Points Temperature",
            0.0,
            1.0,
            0.1,
            label_visibility="collapsed"
        )

        if ppt_impact_container.button('Generate Impact Points'):
            with st.spinner("Generating impact points..."):
                try:
                    # 1. Get the user prompt string
                    user_prompt_for_impact = hlt.generate_prompt(
                        content=content_dict["content"],
                        prompt_name="ppt_impact"
                        # No additional_content needed for ppt_impact
                    )

                    # 2. Instantiate the parser with the new model
                    parser = PydanticOutputParser(pydantic_object=ImpactPoints)

                    # 3. Call the structured generation function
                    structured_result = hlt.generate_structured_content(
                        client=st.session_state.client,
                        system_scope=prompts.SYSTEM_SCOPE,
                        user_prompt=user_prompt_for_impact,
                        pydantic_parser=parser,
                        max_tokens=300, # Adjust as needed
                        temperature=ppt_impact_temperature,
                        max_allowable_tokens=st.session_state.max_allowable_tokens,
                        model=st.session_state.model,
                        package=st.session_state.package
                    )

                    # 4. Store the list of points
                    st.session_state.ppt_impact_response = structured_result.points
                    # ppt_impact_container.success("Impact points generated!")

                except Exception as e:
                    st.session_state.ppt_impact_response = None
                    # ppt_impact_container.error(f"Failed to generate impact points: {e}")


        # --- Updated Display Logic (Editable Text Area) ---
        if st.session_state.ppt_impact_response is not None:
            ppt_impact_container.markdown("Impact Points Result:")
            response_data = st.session_state.ppt_impact_response # Should be a list

            # Prepare string value for the text_area
            if isinstance(response_data, list):
                current_text_value = "\n".join([
                    f"{str(point).strip()}" if str(point).strip().startswith("-") else f"- {str(point).strip()}"
                    for point in response_data if str(point).strip()
                ])
            elif isinstance(response_data, str): # Fallback
                current_text_value = response_data
                # ppt_impact_container.warning("Impact data was a string, expected a list.")
            else: # Handle unexpected
                ppt_impact_container.error(f"Cannot display impact result: Unexpected data type {type(response_data)}")
                current_text_value = ""

            # Use st.text_area for display and editing
            edited_text = ppt_impact_container.text_area(
                label="Impact Points Result:",
                value=current_text_value,
                height=250, # Adjust height
                key="ppt_impact_edit_area", # Unique key
                label_visibility="collapsed"
            )

            # Parse the edited text back into a LIST and update session state
            updated_impact_list = [
                line.strip().lstrip('- ')
                for line in edited_text.split('\n') if line.strip()
            ]
            st.session_state.ppt_impact_response = updated_impact_list

# ------------------------------------------------
# -- PPT:  START FIGURE SELECTION --> 
# ------------------------------------------------

        st.markdown("---")

        # --- New Figure Selection and Caption Section ---
        st.markdown("##### Select Figure and Generate Caption for PowerPoint:")

        figure_select_container = st.container(border=True) # Use border for visual grouping
        ppt_figure_container = st.container(border=True)

        # --- Step 1: List Figure IDs/Descriptions (using LLM) ---
        ppt_figure_container.markdown("##### 1. List Figures Found in Paper")
        if ppt_figure_container.button("List Figures from Text", key="ppt_list_figs"):
            # This reuses the logic from your existing PPT figure section
            with st.spinner("Extracting figure list from paper text..."):
                try:
                    # (Existing logic calling hlt.generate_prompt + hlt.generate_prompt_content
                    #  using the 'figure_list' prompt to populate st.session_state.figure_data)
                    # ... Example call structure ...
                    user_prompt_for_figlist = hlt.generate_prompt(content=content_dict["content"], prompt_name="figure_list")
                    figure_list_raw = hlt.generate_prompt_content(client=st.session_state.client, system_scope=prompts.SYSTEM_SCOPE, prompt=user_prompt_for_figlist, max_tokens=1000, temperature=0.1, max_allowable_tokens=st.session_state.max_allowable_tokens, model=st.session_state.model, package=st.session_state.package)
                    parsed_figures = {}
                    lines = figure_list_raw.strip().split('\n')
                    # (Existing parsing logic for 'Identifier :: Description' format)
                    for line in lines:
                        if line.strip().lower().startswith("table"): continue
                        if ' :: ' in line:
                            parts = line.split(' :: ', 1)
                            identifier = parts[0].strip()
                            description = parts[1].strip()
                            if identifier and description and not identifier.lower().startswith("table"):
                                parsed_figures[identifier] = description
                    st.session_state.figure_data = parsed_figures # Store the dict
                    st.session_state.selected_figure_id = None # Reset selection
                    st.session_state.selected_figure_caption = ""
                    st.session_state.ppt_figure_image_bytes = None # Reset selected image bytes
                    if not st.session_state.figure_data: ppt_figure_container.warning("Could not extract figure list from text.")
                    else: ppt_figure_container.success(f"Found {len(st.session_state.figure_data)} figure references in text.")
                except Exception as e:
                    ppt_figure_container.error(f"Error listing figures: {e}")
                    st.session_state.figure_data = None


        # --- Step 2: Select Figure ID from Dropdown ---
        if st.session_state.figure_data:
            ppt_figure_container.markdown("##### 2. Select Figure ID")
            display_options = ["<Select a Figure ID>"] + [f"{id}: {desc}" for id, desc in st.session_state.figure_data.items()]
            id_lookup = {f"{id}: {desc}": id for id, desc in st.session_state.figure_data.items()}
            current_display_selection = "<Select a Figure ID>"
            # Find current selection display string
            if st.session_state.selected_figure_id and st.session_state.selected_figure_id in st.session_state.figure_data:
                current_display_selection = f"{st.session_state.selected_figure_id}: {st.session_state.figure_data[st.session_state.selected_figure_id]}"

            selected_display_string = ppt_figure_container.selectbox(
                "Choose the Figure ID for the PowerPoint slide:",
                options=display_options,
                index=display_options.index(current_display_selection) if current_display_selection in display_options else 0,
                key="ppt_select_fig_id",
                label_visibility="collapsed"
            )
            # Update selected ID state
            new_selected_id = id_lookup.get(selected_display_string, None)
            if st.session_state.selected_figure_id != new_selected_id:
                st.session_state.selected_figure_id = new_selected_id
                st.session_state.selected_figure_caption = "" # Reset caption
                st.session_state.ppt_figure_image_bytes = None # Reset image selection
                st.rerun() # Rerun if selection changed

        # --- Step 3: Generate Caption for Selected Figure ID ---
        if st.session_state.selected_figure_id: # Only allow caption gen if ID is selected
            ppt_figure_container.markdown("##### 3. Generate Editorial Figure Caption")
            # (This reuses the logic for generating caption based on selected_figure_id
            # and populating selected_figure_caption - ensure button key is unique)
            caption_subcontainer = ppt_figure_container.container() # Separate container
            if caption_subcontainer.button(f"Suggest Caption for {st.session_state.selected_figure_id}", key="gen_ppt_fig_caption"):
                with st.spinner("Generating caption..."):
                    user_prompt_caption = hlt.generate_prompt(content=content_dict["content"], prompt_name="selected_figure_caption", additional_content=st.session_state.selected_figure_id)
                    caption_response = hlt.generate_prompt_content(client=st.session_state.client, system_scope=prompts.SYSTEM_SCOPE, prompt=user_prompt_caption, max_tokens=150, temperature=0.2, max_allowable_tokens=st.session_state.max_allowable_tokens, model=st.session_state.model, package=st.session_state.package)
                    st.session_state.selected_figure_caption = caption_response.strip()
                    st.rerun()

            # Display editable caption
            current_ppt_caption = st.session_state.selected_figure_caption if len(st.session_state.selected_figure_caption) > 0 else ""
            edited_ppt_caption = caption_subcontainer.text_area(
                "Caption:", value=current_ppt_caption, key="edit_ppt_caption", height=100,
                placeholder=f"Enter caption for {st.session_state.selected_figure_id} or generate suggestion..."
            )
            if edited_ppt_caption != current_ppt_caption:
                st.session_state.selected_figure_caption = edited_ppt_caption
                st.rerun()

# ------------------------------------------------
# -- PPT:  START EXPORT -->
# ------------------------------------------------

        st.markdown("---")

        # Add PowerPoint export container at the end
        export_ppt_container = st.container()
        export_ppt_container.markdown("##### Export PowerPoint Presentation with New Content")

        if (st.session_state.title_response is not None and
            st.session_state.objective_response is not None and
            st.session_state.ppt_impact_response is not None and
            st.session_state.approach_response is not None):

            try:
                # Load the PowerPoint template
                ppt_template_file = importlib.resources.files('highlight.data').joinpath('highlight_template.pptx')
                prs = Presentation(ppt_template_file)

                # --- Use the approach list directly ---
                # Ensure approach_response is actually a list
                approach_points = []
                if isinstance(st.session_state.approach_response, list):
                    approach_points = st.session_state.approach_response
                elif st.session_state.approach_response: # Handle case if it's still a string somehow?
                     approach_points = st.session_state.approach_response.split('\n') # Fallback
                else:
                    approach_points = ["Approach not generated."] # Default if missing

                impact_points = [] # Define impact_points similarly
                if isinstance(st.session_state.ppt_impact_response, list): # Assuming ppt_impact is also structured
                    impact_points = st.session_state.ppt_impact_response
                elif st.session_state.ppt_impact_response:
                    impact_points = st.session_state.ppt_impact_response.split('\n')
                else:
                    impact_points = ["Impact points not generated."]

                # --- Process Slides ---
                for slide in prs.slides:

                    # --- Find Placeholders by Name ---
                    impact_ph = get_placeholder(slide, "Text Placeholder 10")
                    approach_ph = get_placeholder(slide, "Text Placeholder 9")
                    picture_ph = get_placeholder(slide, "Picture Placeholder 2")
                    caption_ph = get_placeholder(slide, "Text Placeholder 3")
                    citation_ph = get_placeholder(slide, "Text Placeholder 11")
                    objective_ph = get_placeholder(slide, "Text Placeholder 8")
                    title_ph = get_placeholder(slide, "Title 1")

                    # title
                    tf = title_ph.text_frame
                    tf.clear()
                    tf.text = st.session_state.title_response

                    # objective
                    tf = objective_ph.text_frame
                    tf.clear()
                    tf.text = st.session_state.objective_response
                    
                    # caption
                    tf = caption_ph.text_frame
                    tf.clear()
                    tf.text = st.session_state.selected_figure_caption

                    # citation
                    tf = citation_ph.text_frame
                    tf.clear()
                    tf.text = st.session_state.citation

                    # --- Populate List Placeholders ---
                    if approach_ph and hasattr(approach_ph, 'text_frame') and approach_points:
                        tf = approach_ph.text_frame
                        tf.clear()
                        if len(tf.paragraphs): # Remove potentially empty first paragraph after clear
                             p = tf.paragraphs[0]
                             if not p.text.strip() and len(p.runs) == 0:
                                 p._element.getparent().remove(p._element)
                        # Add points
                        for point_text in approach_points[:3]: # Limit to 3
                            p = tf.add_paragraph()
                            p.text = point_text.strip().lstrip('- ')
                            p.level = 0 # Apply bullet style from template
                            # Optional: Apply specific font overrides if needed
                            p.font.size = Pt(13)
                            p.alignment = PP_ALIGN.LEFT

                    if impact_ph and hasattr(impact_ph, 'text_frame') and impact_points:
                        tf = impact_ph.text_frame
                        tf.clear()
                        if len(tf.paragraphs): # Remove potentially empty first paragraph
                             p = tf.paragraphs[0]
                             if not p.text.strip() and len(p.runs) == 0:
                                 p._element.getparent().remove(p._element)
                        # Add points
                        for point_text in impact_points[:3]: # Limit to 3
                            p = tf.add_paragraph()
                            p.text = point_text.strip().lstrip('- ')
                            p.level = 0 # Apply bullet style from template
                            # Optional: Apply specific font overrides if needed
                            p.font.size = Pt(13)
                            p.alignment = PP_ALIGN.LEFT

                    # --- Populate Picture Placeholder ---
                    if picture_ph:
                         # Check if image bytes exist in the dedicated PPT state variable
                         if st.session_state.get("ppt_figure_image_bytes") is not None:
                            try:
                                # Create BytesIO stream from stored bytes
                                image_stream = io.BytesIO(st.session_state.ppt_figure_image_bytes)
                                image_stream.seek(0)
                                # This replaces the placeholder shape with the picture
                                # May need to adjust size/position afterwards depending on placeholder type
                                picture_ph.insert_picture(image_stream)
                                export_container.info(f"Inserted assigned image into '{picture_ph.name}'.") # Debug message
                            except Exception as pic_e:
                                export_container.warning(f"Could not insert picture into '{picture_ph.name}': {pic_e}")
                         else:
                            export_container.warning(f"Picture placeholder '{picture_ph.name}' found, but no image was assigned.")
                    # --- End Picture Placeholder ---

                # Save the modified presentation to a BytesIO object
                ppt_io = io.BytesIO()
                prs.save(ppt_io)
                ppt_io.seek(0)

                # Provide a download button for the user
                ppt_export_success = export_ppt_container.download_button(
                    label="Export PowerPoint Presentation",
                    data=ppt_io,
                    file_name=f"{st.session_state.base_export_filename}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

                if ppt_export_success:
                    export_ppt_container.success("PowerPoint presentation generated successfully!", icon="✅")

            except Exception as e:
                export_ppt_container.error(f"An error occurred while generating the PowerPoint: {e}", icon="🚨")

        else:
            export_ppt_container.warning("Please generate the title, objective, impact, and citation responses before exporting.", icon="⚠️")
